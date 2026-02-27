#!/usr/bin/env python3
"""
飞书 → Clawdbot 桥接服务
"""

import json
import hashlib
import requests
from flask import Flask, request, jsonify
from functools import lru_cache
import time
import sys
import os
import tempfile
import pandas as pd
import fitz  # PyMuPDF for PDF
import docx  # python-docx for Word
from pptx import Presentation  # python-pptx for PowerPoint

# 禁用输出缓冲
sys.stdout.reconfigure(line_buffering=True)
sys.stderr.reconfigure(line_buffering=True)

# ============ 配置 ============
FEISHU_APP_ID = "cli_a9092df50578dbb3"           # 飞书 App ID
FEISHU_APP_SECRET = "ZPKrfUHhqIL1u8pbcUi0jbtkefKCTXhv"   # 飞书 App Secret
FEISHU_VERIFICATION_TOKEN = "rSxEBMVZbluflS73rVXRndJKV3x2FMXL"  # 事件订阅验证 token

CLAWDBOT_URL = "http://127.0.0.1:18789"  # Clawdbot 地址
CLAWDBOT_TOKEN = "d910603eaf8f32a4e02519cb6a1032741a7efd0cadabd8f9"  # Clawdbot token

# 白名单配置
WHITELIST_FILE = os.path.join(os.path.dirname(__file__), "whitelist_feishu.json")

# 权限配置
PERMISSIONS_FILE = os.path.join(os.path.dirname(__file__), "permissions.json")

# Bot Admin 配置
BOT_ADMIN_URL = "http://127.0.0.1:5002"

# 待审批请求文件
PENDING_FILE = os.path.join(os.path.dirname(__file__), "pending_requests.json")

# 用户会话状态 (用于多轮对话)
user_sessions = {}

# 消息去重缓存 (防止飞书 webhook 重复推送)
from collections import OrderedDict
import threading

class MessageDeduplicator:
    """消息去重器，基于 event_id 去重"""
    def __init__(self, max_size=1000, ttl_seconds=300):
        self.cache = OrderedDict()
        self.max_size = max_size
        self.ttl = ttl_seconds
        self.lock = threading.Lock()
    
    def is_duplicate(self, event_id: str) -> bool:
        """检查是否重复消息，如果不重复则记录"""
        if not event_id:
            return False
        
        now = time.time()
        with self.lock:
            # 清理过期条目
            expired = [k for k, v in self.cache.items() if now - v > self.ttl]
            for k in expired:
                del self.cache[k]
            
            # 检查是否存在
            if event_id in self.cache:
                return True
            
            # 添加新条目
            self.cache[event_id] = now
            
            # 限制大小
            while len(self.cache) > self.max_size:
                self.cache.popitem(last=False)
            
            return False

message_dedup = MessageDeduplicator()
# ==============================

app = Flask(__name__)


# ============ 白名单功能 ============
whitelist_cache = {"data": None, "mtime": 0}


def load_whitelist():
    """加载白名单，支持热更新（文件修改后自动重新加载）"""
    try:
        if not os.path.exists(WHITELIST_FILE):
            return {"enabled": False, "users": []}
        
        mtime = os.path.getmtime(WHITELIST_FILE)
        if whitelist_cache["data"] and whitelist_cache["mtime"] == mtime:
            return whitelist_cache["data"]
        
        with open(WHITELIST_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        
        whitelist_cache["data"] = data
        whitelist_cache["mtime"] = mtime
        print(f"白名单已加载: {len(data.get('users', []))} 个用户, enabled={data.get('enabled', False)}")
        return data
    except Exception as e:
        print(f"加载白名单失败: {e}")
        return {"enabled": False, "users": []}


def is_user_allowed(sender_info: dict, sender_type: str = "user") -> bool:
    """检查用户是否在白名单中
    
    Args:
        sender_info: 发送者信息 (包含 user_id, open_id 等)
        sender_type: 发送者类型 ("user" 或 "app")
    """
    whitelist = load_whitelist()
    
    # 白名单未启用，允许所有人
    if not whitelist.get("enabled", False):
        return True
    
    users = whitelist.get("users", [])
    bots = whitelist.get("bots", [])  # 允许的 bot open_id 列表
    
    # 白名单为空，允许所有人
    if not users and not bots:
        return True
    
    # 检查 user_id, open_id, union_id 是否在白名单中
    user_id = sender_info.get("user_id", "")
    open_id = sender_info.get("open_id", "")
    union_id = sender_info.get("union_id", "")
    
    # 如果是 bot 消息，检查 bots 白名单
    if sender_type == "app":
        return open_id in bots
    
    return user_id in users or open_id in users or union_id in users


def is_verification_code(text: str) -> bool:
    """检查是否是验证码格式 (6位大写字母+数字)"""
    import re
    text = text.strip().upper()
    return bool(re.match(r'^[A-Z0-9]{6}$', text))


def verify_code_with_admin(code: str, open_id: str, chat_id: str = None) -> dict:
    """调用 bot-admin 验证验证码"""
    try:
        resp = requests.post(
            f"{BOT_ADMIN_URL}/verify",
            json={"code": code.upper(), "open_id": open_id, "chat_id": chat_id},
            timeout=10
        )
        return resp.json()
    except Exception as e:
        print(f"调用 bot-admin 失败: {e}")
        return {"found": False, "message": "验证服务暂时不可用"}


# ============ 权限管理功能 ============
permissions_cache = {"data": None, "mtime": 0}


def load_permissions():
    """加载权限配置，支持热更新"""
    try:
        if not os.path.exists(PERMISSIONS_FILE):
            return {"roles": {}, "features": {}, "users": {}}
        
        mtime = os.path.getmtime(PERMISSIONS_FILE)
        if permissions_cache["data"] and permissions_cache["mtime"] == mtime:
            return permissions_cache["data"]
        
        with open(PERMISSIONS_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        
        permissions_cache["data"] = data
        permissions_cache["mtime"] = mtime
        print(f"权限配置已加载: {len(data.get('users', {}))} 个用户")
        return data
    except Exception as e:
        print(f"加载权限配置失败: {e}")
        return {"roles": {}, "features": {}, "users": {}}


def get_user_permissions(open_id: str) -> dict:
    """获取用户的权限信息"""
    perms = load_permissions()
    
    # 查找用户
    user_info = perms.get("users", {}).get(open_id, {})
    if not user_info:
        # 默认给 viewer 角色
        user_info = {"name": "unknown", "role": "viewer"}
    
    role_name = user_info.get("role", "viewer")
    role_info = perms.get("roles", {}).get(role_name, {})
    features = role_info.get("features", [])
    
    # 处理通配符
    if "*" in features:
        features = list(perms.get("features", {}).keys())
    
    return {
        "name": user_info.get("name", "unknown"),
        "role": role_name,
        "role_description": role_info.get("description", ""),
        "features": features
    }


def has_permission(open_id: str, feature: str) -> bool:
    """检查用户是否有某项功能权限"""
    user_perms = get_user_permissions(open_id)
    features = user_perms.get("features", [])
    return "*" in features or feature in features


# ============ 权限申请功能 ============
def load_pending():
    """加载待审批请求"""
    try:
        if os.path.exists(PENDING_FILE):
            with open(PENDING_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception as e:
        print(f"加载待审批请求失败: {e}")
    return {"pending": {}}


def save_pending(data):
    """保存待审批请求"""
    with open(PENDING_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def get_admin_chat_ids():
    """获取所有 admin 用户的 chat_id (需要他们先私聊过机器人)"""
    perms = load_permissions()
    admin_ids = []
    for open_id, info in perms.get("users", {}).items():
        if info.get("role") == "admin" and info.get("chat_id"):
            admin_ids.append(info["chat_id"])
    return admin_ids


def save_user_chat_id(open_id, chat_id):
    """保存用户的 chat_id (用于后续通知)"""
    perms = load_permissions()
    if open_id in perms.get("users", {}):
        perms["users"][open_id]["chat_id"] = chat_id
        with open(PERMISSIONS_FILE, "w", encoding="utf-8") as f:
            json.dump(perms, f, ensure_ascii=False, indent=2)


def add_to_whitelist(open_id, name):
    """添加用户到白名单"""
    whitelist = load_whitelist()
    if open_id not in whitelist.get("users", []):
        whitelist.setdefault("users", []).append(open_id)
        note = whitelist.get("note", "")
        whitelist["note"] = f"{note} / {name}" if note else name
        with open(WHITELIST_FILE, "w", encoding="utf-8") as f:
            json.dump(whitelist, f, ensure_ascii=False, indent=2)
        # 清除缓存
        whitelist_cache["data"] = None


def add_to_permissions(open_id, name, role):
    """添加用户到权限系统"""
    perms = load_permissions()
    if "users" not in perms:
        perms["users"] = {}
    perms["users"][open_id] = {"name": name, "role": role}
    with open(PERMISSIONS_FILE, "w", encoding="utf-8") as f:
        json.dump(perms, f, ensure_ascii=False, indent=2)
    # 清除缓存
    permissions_cache["data"] = None


def send_approval_card(admin_chat_id, applicant_name, applicant_open_id, requested_role, chat_id):
    """发送审批卡片给 Admin"""
    token = get_tenant_access_token()
    if not token:
        return False
    
    # 简化的卡片格式 - 每个按钮代表一个角色
    card = {
        "config": {"wide_screen_mode": True},
        "header": {
            "title": {"tag": "plain_text", "content": "📋 新权限申请"},
            "template": "blue"
        },
        "elements": [
            {
                "tag": "div",
                "text": {"tag": "lark_md", "content": f"**用户:** {applicant_name}\n**申请角色:** {requested_role}"}
            },
            {"tag": "hr"},
            {
                "tag": "div",
                "text": {"tag": "lark_md", "content": "点击按钮批准（使用对应角色）或拒绝："}
            },
            {
                "tag": "action",
                "actions": [
                    {
                        "tag": "button",
                        "text": {"tag": "plain_text", "content": "✅ viewer"},
                        "type": "default",
                        "value": json.dumps({"action": "approve", "role": "viewer", "open_id": applicant_open_id, "name": applicant_name, "chat_id": chat_id})
                    },
                    {
                        "tag": "button",
                        "text": {"tag": "plain_text", "content": "✅ user"},
                        "type": "default",
                        "value": json.dumps({"action": "approve", "role": "user", "open_id": applicant_open_id, "name": applicant_name, "chat_id": chat_id})
                    },
                    {
                        "tag": "button",
                        "text": {"tag": "plain_text", "content": "✅ power_user"},
                        "type": "primary",
                        "value": json.dumps({"action": "approve", "role": "power_user", "open_id": applicant_open_id, "name": applicant_name, "chat_id": chat_id})
                    }
                ]
            },
            {
                "tag": "action",
                "actions": [
                    {
                        "tag": "button",
                        "text": {"tag": "plain_text", "content": "✅ admin"},
                        "type": "primary",
                        "value": json.dumps({"action": "approve", "role": "admin", "open_id": applicant_open_id, "name": applicant_name, "chat_id": chat_id})
                    },
                    {
                        "tag": "button",
                        "text": {"tag": "plain_text", "content": "❌ 拒绝"},
                        "type": "danger",
                        "value": json.dumps({"action": "reject", "open_id": applicant_open_id, "name": applicant_name, "chat_id": chat_id})
                    }
                ]
            }
        ]
    }
    
    url = "https://open.feishu.cn/open-apis/im/v1/messages"
    headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
    payload = {
        "receive_id": admin_chat_id,
        "msg_type": "interactive",
        "content": json.dumps(card)
    }
    
    resp = requests.post(url, headers=headers, json=payload, params={"receive_id_type": "chat_id"})
    print(f"发送审批卡片结果: {resp.json()}")
    return resp.json()


ROLE_DESCRIPTIONS = {
    "viewer": "只能搜索和对话",
    "user": "可以搜索和读取文件", 
    "power_user": "可以读写文件（无命令执行）",
    "admin": "全部权限"
}


def handle_permission_request(text, open_id, chat_id, sender_name):
    """处理权限申请流程"""
    session = user_sessions.get(open_id, {})
    
    # 检查是否已有待审批的申请
    pending = load_pending()
    if open_id in pending.get("pending", {}):
        existing = pending["pending"][open_id]
        return f"你已有一个待审批的申请（角色: {existing.get('requested_role', 'unknown')}）\n\n请等待管理员审批，无需重复申请。"
    
    # 检查是否在申请流程中
    if session.get("state") == "selecting_role":
        # 用户正在选择角色
        role_map = {"1": "viewer", "2": "user", "3": "power_user"}
        selected_role = role_map.get(text.strip())
        
        if selected_role:
            # 保存申请
            pending["pending"][open_id] = {
                "name": sender_name,
                "requested_role": selected_role,
                "chat_id": chat_id,
                "created_at": time.time()
            }
            save_pending(pending)
            
            # 清除会话状态
            user_sessions.pop(open_id, None)
            
            # 通知所有 Admin
            admin_chat_ids = get_admin_chat_ids()
            if admin_chat_ids:
                for admin_chat_id in admin_chat_ids:
                    send_approval_card(admin_chat_id, sender_name, open_id, selected_role, chat_id)
                return f"✅ 申请已提交！\n\n你申请的角色: {selected_role} ({ROLE_DESCRIPTIONS[selected_role]})\n\n请等待管理员审批。"
            else:
                return "✅ 申请已提交！\n\n⚠️ 暂无在线管理员，请联系管理员手动处理。"
        else:
            return "请回复数字 1、2 或 3 选择角色：\n\n1. viewer - 只能搜索和对话\n2. user - 可以搜索和读取文件\n3. power_user - 可以读写文件"
    
    # 新申请
    if text.strip() in ["申请权限", "/request", "申请使用", "申请"]:
        user_sessions[open_id] = {"state": "selecting_role", "created_at": time.time()}
        return "请选择你想申请的角色：\n\n1. viewer - 只能搜索和对话\n2. user - 可以搜索和读取文件\n3. power_user - 可以读写文件\n\n回复数字即可（如：1）"
    
    return None  # 不是权限申请相关的消息
# =====================================

# 缓存 tenant_access_token
token_cache = {"token": None, "expire": 0}


def get_tenant_access_token():
    """获取飞书 tenant_access_token"""
    now = time.time()
    if token_cache["token"] and token_cache["expire"] > now:
        return token_cache["token"]
    
    url = "https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal"
    resp = requests.post(url, json={
        "app_id": FEISHU_APP_ID,
        "app_secret": FEISHU_APP_SECRET
    })
    data = resp.json()
    
    if data.get("code") == 0:
        token_cache["token"] = data["tenant_access_token"]
        token_cache["expire"] = now + data.get("expire", 7200) - 60
        return token_cache["token"]
    else:
        print(f"获取 token 失败: {data}")
        return None


def send_feishu_message(chat_id: str, text: str):
    """发送消息到飞书（通过 chat_id）"""
    token = get_tenant_access_token()
    if not token:
        return False
    
    url = "https://open.feishu.cn/open-apis/im/v1/messages"
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json"
    }
    payload = {
        "receive_id": chat_id,
        "msg_type": "text",
        "content": json.dumps({"text": text})
    }
    
    resp = requests.post(url, headers=headers, json=payload, params={"receive_id_type": "chat_id"})
    result = resp.json()
    print(f"发送消息结果: {result}")
    return result


def send_private_message(open_id: str, text: str):
    """发送私聊消息给用户（通过 open_id）"""
    token = get_tenant_access_token()
    if not token:
        return False
    
    url = "https://open.feishu.cn/open-apis/im/v1/messages"
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json"
    }
    payload = {
        "receive_id": open_id,
        "msg_type": "text",
        "content": json.dumps({"text": text})
    }
    
    resp = requests.post(url, headers=headers, json=payload, params={"receive_id_type": "open_id"})
    result = resp.json()
    print(f"发送私聊消息结果: {result}")
    return result


def upload_file_to_feishu(file_path: str, file_name: str = None) -> str:
    """上传文件到飞书，返回 file_key"""
    token = get_tenant_access_token()
    if not token:
        print("获取 token 失败")
        return None
    
    if not os.path.exists(file_path):
        print(f"文件不存在: {file_path}")
        return None
    
    if not file_name:
        file_name = os.path.basename(file_path)
    
    url = "https://open.feishu.cn/open-apis/im/v1/files"
    headers = {"Authorization": f"Bearer {token}"}
    
    # 根据扩展名确定 file_type
    ext = os.path.splitext(file_name)[1].lower()
    if ext in [".xls", ".xlsx"]:
        file_type = "xls"
    elif ext in [".pdf"]:
        file_type = "pdf"
    elif ext in [".doc", ".docx"]:
        file_type = "doc"
    elif ext in [".ppt", ".pptx"]:
        file_type = "ppt"
    elif ext in [".mp4", ".mov", ".avi", ".mkv"]:
        file_type = "mp4"
    elif ext in [".mp3", ".wav", ".ogg", ".m4a"]:
        file_type = "opus"
    else:
        file_type = "stream"  # 通用二进制
    
    try:
        with open(file_path, "rb") as f:
            files = {
                "file": (file_name, f),
            }
            data = {
                "file_type": file_type,
                "file_name": file_name
            }
            resp = requests.post(url, headers=headers, files=files, data=data, timeout=60)
            result = resp.json()
            print(f"上传文件结果: {result}")
            
            if result.get("code") == 0:
                return result.get("data", {}).get("file_key")
            else:
                print(f"上传失败: {result}")
                return None
    except Exception as e:
        print(f"上传文件异常: {e}")
        return None


def send_feishu_file(receive_id: str, file_key: str, file_name: str = "file", file_type: str = "file", receive_id_type: str = "chat_id"):
    """发送文件消息到飞书
    
    Args:
        receive_id: 接收者ID (chat_id 或 open_id)
        file_key: 文件key
        file_name: 文件名
        file_type: 文件类型 (mp4/opus 用 media, 其他用 file)
        receive_id_type: 接收者类型 (chat_id 或 open_id)
    """
    token = get_tenant_access_token()
    if not token:
        return False
    
    url = "https://open.feishu.cn/open-apis/im/v1/messages"
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json"
    }
    
    # 视频和音频用 media 类型发送
    if file_type in ["mp4", "opus"]:
        msg_type = "media"
    else:
        msg_type = "file"
    
    payload = {
        "receive_id": receive_id,
        "msg_type": msg_type,
        "content": json.dumps({"file_key": file_key})
    }
    
    resp = requests.post(url, headers=headers, json=payload, params={"receive_id_type": receive_id_type})
    result = resp.json()
    print(f"发送{msg_type}结果: {result}")
    return result


def send_file_to_chat(chat_id: str = None, file_path: str = None, file_name: str = None, caption: str = None, open_id: str = None):
    """上传并发送文件到飞书聊天（完整流程）
    
    Args:
        chat_id: 群聊ID (与 open_id 二选一)
        file_path: 文件路径
        file_name: 文件名
        caption: 说明文字
        open_id: 用户ID，用于私聊发送 (与 chat_id 二选一)
    """
    if not file_name:
        file_name = os.path.basename(file_path)
    
    # 确定接收者
    if open_id:
        receive_id = open_id
        receive_id_type = "open_id"
    else:
        receive_id = chat_id
        receive_id_type = "chat_id"
    
    # 确定文件类型
    ext = os.path.splitext(file_name)[1].lower()
    if ext in [".mp4", ".mov", ".avi", ".mkv"]:
        file_type = "mp4"
    elif ext in [".mp3", ".wav", ".ogg", ".m4a"]:
        file_type = "opus"
    else:
        file_type = "file"
    
    # 1. 上传文件
    file_key = upload_file_to_feishu(file_path, file_name)
    if not file_key:
        return {"success": False, "error": "文件上传失败"}
    
    # 2. 发送说明文字（如果有）
    if caption:
        if open_id:
            send_private_message(open_id, caption)
        else:
            send_feishu_message(chat_id, caption)
    
    # 3. 发送文件/媒体
    result = send_feishu_file(receive_id, file_key, file_name, file_type, receive_id_type)
    
    if result.get("code") == 0:
        return {"success": True, "file_key": file_key, "result": result}
    else:
        return {"success": False, "error": result.get("msg", "发送失败")}


# ============ 获取群聊历史消息 ============
def get_recent_messages(chat_id: str, limit: int = 20) -> list:
    """获取群聊最近的消息"""
    token = get_tenant_access_token()
    if not token:
        return []
    
    url = f"https://open.feishu.cn/open-apis/im/v1/messages"
    headers = {"Authorization": f"Bearer {token}"}
    params = {
        "container_id_type": "chat",
        "container_id": chat_id,
        "page_size": limit,
        "sort_type": "ByCreateTimeDesc"  # 最新的在前
    }
    
    try:
        resp = requests.get(url, headers=headers, params=params, timeout=30)
        result = resp.json()
        if result.get("code") == 0:
            items = result.get("data", {}).get("items", [])
            print(f"获取到 {len(items)} 条历史消息")
            return items
        else:
            print(f"获取历史消息失败: {result}")
            return []
    except Exception as e:
        print(f"获取历史消息异常: {e}")
        return []


def find_recent_file(chat_id: str, max_messages: int = 20) -> dict:
    """从最近消息中找到文件消息"""
    messages = get_recent_messages(chat_id, max_messages)
    
    for msg in messages:
        msg_type = msg.get("msg_type")
        if msg_type == "file":
            content = json.loads(msg.get("body", {}).get("content", "{}"))
            file_key = content.get("file_key", "")
            file_name = content.get("file_name", "")
            message_id = msg.get("message_id", "")
            
            if file_key and file_name:
                print(f"找到文件: {file_name}, message_id: {message_id}")
                return {
                    "found": True,
                    "file_key": file_key,
                    "file_name": file_name,
                    "message_id": message_id
                }
    
    return {"found": False}


def is_file_request(text: str) -> bool:
    """检测用户是否在请求读取文件"""
    keywords = [
        "读文件", "看文件", "分析文件", "读一下文件", "看一下文件",
        "读pdf", "看pdf", "分析pdf",
        "读word", "看word", "分析word",
        "读excel", "看excel", "分析excel",
        "读ppt", "看ppt", "分析ppt",
        "读表格", "看表格", "分析表格",
        "读文档", "看文档", "分析文档",
        "上面的文件", "刚才的文件", "那个文件",
        "帮我读", "帮我看", "帮我分析",
        "read file", "read the file", "analyze file",
        "这个文件", "处理文件", "解析文件"
    ]
    text_lower = text.lower()
    return any(kw in text_lower for kw in keywords)


# ============ 文件处理功能 ============
def download_feishu_file(message_id: str, file_key: str, file_name: str) -> str:
    """从飞书下载文件，返回临时文件路径"""
    token = get_tenant_access_token()
    if not token:
        return None
    
    url = f"https://open.feishu.cn/open-apis/im/v1/messages/{message_id}/resources/{file_key}"
    headers = {"Authorization": f"Bearer {token}"}
    params = {"type": "file"}
    
    try:
        resp = requests.get(url, headers=headers, params=params, timeout=60)
        if resp.status_code == 200:
            # 保存到临时文件
            suffix = os.path.splitext(file_name)[1] or ".xlsx"
            fd, temp_path = tempfile.mkstemp(suffix=suffix)
            with os.fdopen(fd, 'wb') as f:
                f.write(resp.content)
            print(f"文件下载成功: {temp_path} ({len(resp.content)} bytes)")
            return temp_path
        else:
            print(f"文件下载失败: {resp.status_code} - {resp.text}")
            return None
    except Exception as e:
        print(f"文件下载异常: {e}")
        return None


def download_feishu_image(message_id: str, image_key: str) -> str:
    """从飞书下载图片，返回临时文件路径"""
    token = get_tenant_access_token()
    if not token:
        return None
    
    url = f"https://open.feishu.cn/open-apis/im/v1/messages/{message_id}/resources/{image_key}"
    headers = {"Authorization": f"Bearer {token}"}
    params = {"type": "image"}
    
    try:
        resp = requests.get(url, headers=headers, params=params, timeout=60)
        if resp.status_code == 200:
            # 检测图片类型
            content_type = resp.headers.get("Content-Type", "image/png")
            if "jpeg" in content_type or "jpg" in content_type:
                suffix = ".jpg"
            elif "gif" in content_type:
                suffix = ".gif"
            elif "webp" in content_type:
                suffix = ".webp"
            else:
                suffix = ".png"
            
            fd, temp_path = tempfile.mkstemp(suffix=suffix)
            with os.fdopen(fd, 'wb') as f:
                f.write(resp.content)
            print(f"图片下载成功: {temp_path} ({len(resp.content)} bytes)")
            return temp_path
        else:
            print(f"图片下载失败: {resp.status_code} - {resp.text}")
            return None
    except Exception as e:
        print(f"图片下载异常: {e}")
        return None


def parse_excel_file(file_path: str, max_rows: int = 100) -> str:
    """解析 Excel 文件，返回文本内容"""
    try:
        # 根据扩展名选择引擎
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".xls":
            engine = "xlrd"
        else:
            engine = "openpyxl"
        
        # 读取所有 sheet
        excel_file = pd.ExcelFile(file_path, engine=engine)
        sheet_names = excel_file.sheet_names
        
        result_parts = []
        total_rows = 0
        
        for sheet_name in sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            
            if df.empty:
                continue
            
            # 限制行数
            remaining_rows = max_rows - total_rows
            if remaining_rows <= 0:
                result_parts.append(f"\n[已达到 {max_rows} 行限制，后续内容省略...]")
                break
            
            if len(df) > remaining_rows:
                df = df.head(remaining_rows)
                truncated = True
            else:
                truncated = False
            
            total_rows += len(df)
            
            # 格式化输出
            if len(sheet_names) > 1:
                result_parts.append(f"\n【Sheet: {sheet_name}】")
            
            # 转为文本表格
            result_parts.append(df.to_string(index=False, max_colwidth=50))
            
            if truncated:
                result_parts.append(f"\n[Sheet {sheet_name} 数据已截断...]")
        
        if not result_parts:
            return "[Excel 文件为空]"
        
        return "\n".join(result_parts)
    
    except Exception as e:
        print(f"Excel 解析失败: {e}")
        import traceback
        traceback.print_exc()
        return f"[Excel 解析失败: {str(e)}]"
    finally:
        # 清理临时文件
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except:
            pass


def is_excel_file(file_name: str) -> bool:
    """检查是否是 Excel 文件"""
    ext = os.path.splitext(file_name)[1].lower()
    return ext in [".xls", ".xlsx"]


def is_pdf_file(file_name: str) -> bool:
    """检查是否是 PDF 文件"""
    ext = os.path.splitext(file_name)[1].lower()
    return ext == ".pdf"


def is_word_file(file_name: str) -> bool:
    """检查是否是 Word 文件"""
    ext = os.path.splitext(file_name)[1].lower()
    return ext in [".doc", ".docx"]


def is_ppt_file(file_name: str) -> bool:
    """检查是否是 PowerPoint 文件"""
    ext = os.path.splitext(file_name)[1].lower()
    return ext in [".ppt", ".pptx"]


def is_md_file(file_name: str) -> bool:
    """检查是否是 Markdown 文件"""
    ext = os.path.splitext(file_name)[1].lower()
    return ext in [".md", ".markdown"]


def is_text_file(file_name: str) -> bool:
    """检查是否是纯文本文件"""
    ext = os.path.splitext(file_name)[1].lower()
    return ext in [".txt", ".log", ".json", ".yaml", ".yml", ".csv", ".xml", ".html", ".css", ".js", ".py", ".sh"]


def is_video_file(file_name: str) -> bool:
    """检查是否是视频文件"""
    ext = os.path.splitext(file_name)[1].lower()
    return ext in [".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v"]


def is_audio_file(file_name: str) -> bool:
    """检查是否是音频文件"""
    ext = os.path.splitext(file_name)[1].lower()
    return ext in [".mp3", ".wav", ".ogg", ".opus", ".m4a", ".aac", ".flac"]


def is_supported_file(file_name: str) -> bool:
    """检查是否是支持的文件类型"""
    return (is_excel_file(file_name) or is_pdf_file(file_name) or is_word_file(file_name) or 
            is_ppt_file(file_name) or is_md_file(file_name) or is_text_file(file_name) or
            is_video_file(file_name) or is_audio_file(file_name))


def parse_pdf_file(file_path: str, max_chars: int = 50000) -> str:
    """解析 PDF 文件，返回文本内容"""
    try:
        doc = fitz.open(file_path)
        text_parts = []
        total_chars = 0
        
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text()
            
            if total_chars + len(page_text) > max_chars:
                remaining = max_chars - total_chars
                if remaining > 0:
                    text_parts.append(f"\n【第 {page_num} 页】\n{page_text[:remaining]}")
                text_parts.append(f"\n[已达到 {max_chars} 字符限制，共 {len(doc)} 页，后续内容省略...]")
                break
            
            text_parts.append(f"\n【第 {page_num} 页】\n{page_text}")
            total_chars += len(page_text)
        
        doc.close()
        
        if not text_parts:
            return "[PDF 文件为空或无法提取文本]"
        
        return "".join(text_parts).strip()
    
    except Exception as e:
        print(f"PDF 解析失败: {e}")
        import traceback
        traceback.print_exc()
        return f"[PDF 解析失败: {str(e)}]"
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except:
            pass


def parse_word_file(file_path: str, max_chars: int = 50000) -> str:
    """解析 Word 文件，返回文本内容"""
    try:
        doc = docx.Document(file_path)
        text_parts = []
        total_chars = 0
        
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if not para_text:
                continue
            
            if total_chars + len(para_text) > max_chars:
                remaining = max_chars - total_chars
                if remaining > 0:
                    text_parts.append(para_text[:remaining])
                text_parts.append(f"\n[已达到 {max_chars} 字符限制，后续内容省略...]")
                break
            
            text_parts.append(para_text)
            total_chars += len(para_text)
        
        # 也尝试提取表格内容
        for table in doc.tables:
            if total_chars > max_chars:
                break
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                table_text.append(row_text)
                total_chars += len(row_text)
            if table_text:
                text_parts.append("\n【表格】\n" + "\n".join(table_text))
        
        if not text_parts:
            return "[Word 文件为空]"
        
        return "\n".join(text_parts).strip()
    
    except Exception as e:
        print(f"Word 解析失败: {e}")
        import traceback
        traceback.print_exc()
        return f"[Word 解析失败: {str(e)}]"
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except:
            pass


def parse_ppt_file(file_path: str, max_chars: int = 50000) -> str:
    """解析 PowerPoint 文件，返回文本内容"""
    try:
        prs = Presentation(file_path)
        text_parts = []
        total_chars = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                
                # 处理表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)
            
            if slide_texts:
                slide_content = "\n".join(slide_texts)
                
                if total_chars + len(slide_content) > max_chars:
                    remaining = max_chars - total_chars
                    if remaining > 0:
                        text_parts.append(f"\n【第 {slide_num} 页】\n{slide_content[:remaining]}")
                    text_parts.append(f"\n[已达到 {max_chars} 字符限制，共 {len(prs.slides)} 页，后续内容省略...]")
                    break
                
                text_parts.append(f"\n【第 {slide_num} 页】\n{slide_content}")
                total_chars += len(slide_content)
        
        if not text_parts:
            return "[PPT 文件为空或无法提取文本]"
        
        return "".join(text_parts).strip()
    
    except Exception as e:
        print(f"PPT 解析失败: {e}")
        import traceback
        traceback.print_exc()
        return f"[PPT 解析失败: {str(e)}]"
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except:
            pass


def parse_md_file(file_path: str, max_chars: int = 50000) -> str:
    """解析 Markdown 文件，返回文本内容"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n[已达到 {max_chars} 字符限制，后续内容省略...]"
        
        if not content.strip():
            return "[Markdown 文件为空]"
        
        return content.strip()
    
    except UnicodeDecodeError:
        # 尝试其他编码
        try:
            with open(file_path, 'r', encoding='gbk') as f:
                content = f.read()
            if len(content) > max_chars:
                content = content[:max_chars] + f"\n\n[已达到 {max_chars} 字符限制，后续内容省略...]"
            return content.strip()
        except:
            return "[Markdown 文件编码无法识别]"
    except Exception as e:
        print(f"Markdown 解析失败: {e}")
        import traceback
        traceback.print_exc()
        return f"[Markdown 解析失败: {str(e)}]"
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except:
            pass


def parse_text_file(file_path: str, max_chars: int = 50000) -> str:
    """解析纯文本文件，返回内容"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n[已达到 {max_chars} 字符限制，后续内容省略...]"
        
        if not content.strip():
            return "[文本文件为空]"
        
        return content.strip()
    
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='gbk') as f:
                content = f.read()
            if len(content) > max_chars:
                content = content[:max_chars] + f"\n\n[已达到 {max_chars} 字符限制，后续内容省略...]"
            return content.strip()
        except:
            return "[文本文件编码无法识别]"
    except Exception as e:
        print(f"文本文件解析失败: {e}")
        import traceback
        traceback.print_exc()
        return f"[文本文件解析失败: {str(e)}]"
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except:
            pass


def parse_video_file(file_path: str, file_name: str) -> str:
    """处理视频文件 - 返回文件信息（视频内容无法直接解析为文字）"""
    try:
        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        size_mb = file_size / (1024 * 1024)
        ext = os.path.splitext(file_name)[1].lower()
        return f"[视频文件]\n文件名: {file_name}\n格式: {ext}\n大小: {size_mb:.2f} MB\n\n（视频内容无法直接解析，请描述您需要对这个视频做什么操作）"
    except Exception as e:
        return f"[视频文件: {file_name}，获取信息失败: {str(e)}]"
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except:
            pass


def parse_audio_file(file_path: str, file_name: str) -> str:
    """处理音频文件 - 返回文件信息（音频内容无法直接解析为文字）"""
    try:
        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        size_mb = file_size / (1024 * 1024)
        ext = os.path.splitext(file_name)[1].lower()
        return f"[音频文件]\n文件名: {file_name}\n格式: {ext}\n大小: {size_mb:.2f} MB\n\n（音频内容无法直接解析，请描述您需要对这个音频做什么操作）"
    except Exception as e:
        return f"[音频文件: {file_name}，获取信息失败: {str(e)}]"
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except:
            pass


def parse_document_file(file_path: str, file_name: str) -> tuple:
    """解析文档文件，返回 (内容, 文件类型描述)"""
    if is_excel_file(file_name):
        return parse_excel_file(file_path), "Excel 表格"
    elif is_pdf_file(file_name):
        return parse_pdf_file(file_path), "PDF 文档"
    elif is_word_file(file_name):
        return parse_word_file(file_path), "Word 文档"
    elif is_ppt_file(file_name):
        return parse_ppt_file(file_path), "PPT 演示文稿"
    elif is_md_file(file_name):
        return parse_md_file(file_path), "Markdown 文档"
    elif is_text_file(file_name):
        return parse_text_file(file_path), "文本文件"
    elif is_video_file(file_name):
        return parse_video_file(file_path, file_name), "视频文件"
    elif is_audio_file(file_name):
        return parse_audio_file(file_path, file_name), "音频文件"
    else:
        return None, None
# ============================================


# 角色到 Agent 的映射
ROLE_TO_AGENT = {
    "admin": "clawdbot:main",        # 全部工具
    "ecommerce_ops": "clawdbot:main", # 电商运营管理员 - 可发货、查订单
    "power_user": "clawdbot:power-user",  # 大部分工具，无 exec
    "user": "clawdbot:user",          # 只读 + 搜索
    "viewer": "clawdbot:viewer"       # 仅搜索，无文件访问
}


def ask_clawdbot(message: str, user_id: str, open_id: str = None, chat_id: str = None) -> str:
    """调用 Clawdbot API"""
    try:
        # 获取用户权限
        user_perms = get_user_permissions(open_id) if open_id else {}
        role = user_perms.get("role", "viewer")
        user_name = user_perms.get("name", "unknown")
        
        # 根据角色选择 Agent (硬性权限控制)
        agent = ROLE_TO_AGENT.get(role, "clawdbot:viewer")
        
        # 添加用户信息和聊天上下文到消息
        user_context = f"[飞书消息 | 用户: {user_name} | 角色: {role} | chat_id: {chat_id}]\n\n"
        full_message = user_context + message
        
        print(f"用户 {user_name} (角色: {role}) -> Agent: {agent}")
        
        resp = requests.post(
            f"{CLAWDBOT_URL}/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {CLAWDBOT_TOKEN}",
                "Content-Type": "application/json"
            },
            json={
                "model": agent,  # 根据角色选择不同的 agent
                "user": f"feishu:{user_id}",  # 保持会话
                "messages": [{"role": "user", "content": full_message}]
            },
            timeout=300
        )
        data = resp.json()
        return data["choices"][0]["message"]["content"]
    except Exception as e:
        error_str = str(e)
        print(f"Clawdbot 调用失败: {error_str}")
        
        # 对常见错误返回友好提示
        if "429" in error_str or "rate_limit" in error_str:
            return "系统繁忙，请稍后再试 🙏"
        elif "timeout" in error_str.lower() or "timed out" in error_str.lower():
            return "响应超时，请稍后再试"
        elif "500" in error_str or "502" in error_str or "503" in error_str:
            return "服务暂时不可用，请稍后再试"
        else:
            return "抱歉，处理请求时出现问题，请稍后再试"


def ask_clawdbot_with_image(image_path: str, user_id: str, open_id: str = None, chat_id: str = None, prompt: str = None) -> str:
    """调用 Clawdbot API 并附带图片"""
    import shutil
    
    try:
        # 获取用户权限
        user_perms = get_user_permissions(open_id) if open_id else {}
        role = user_perms.get("role", "viewer")
        user_name = user_perms.get("name", "unknown")
        
        # 根据角色选择 Agent
        agent = ROLE_TO_AGENT.get(role, "clawdbot:viewer")
        
        # 保存图片到固定位置，让 Clawdbot 可以读取
        import uuid
        ext = os.path.splitext(image_path)[1].lower() or ".jpg"
        image_id = str(uuid.uuid4())[:8]
        saved_path = f"/home/juhe0092/clawd/temp/feishu_image_{image_id}{ext}"
        
        # 确保目录存在
        os.makedirs(os.path.dirname(saved_path), exist_ok=True)
        shutil.copy(image_path, saved_path)
        print(f"图片已保存到: {saved_path}")
        
        # 构造消息，告诉 Clawdbot 图片位置
        text_content = prompt or "请分析这张图片"
        user_context = f"[飞书消息 | 用户: {user_name} | 角色: {role} | chat_id: {chat_id}]\n\n"
        
        # 明确告诉 Clawdbot 读取图片文件
        full_message = f"{user_context}用户发送了一张图片，保存在: {saved_path}\n\n用户说: {text_content}\n\n请先用 Read 工具读取这张图片，然后根据用户的要求处理。"
        
        print(f"用户 {user_name} (角色: {role}) 发送图片 -> Agent: {agent}")
        
        resp = requests.post(
            f"{CLAWDBOT_URL}/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {CLAWDBOT_TOKEN}",
                "Content-Type": "application/json"
            },
            json={
                "model": agent,
                "user": f"feishu:{user_id}",
                "messages": [{"role": "user", "content": full_message}]
            },
            timeout=300
        )
        data = resp.json()
        return data["choices"][0]["message"]["content"]
    except Exception as e:
        error_str = str(e)
        print(f"Clawdbot 图片调用失败: {error_str}")
        
        if "429" in error_str or "rate_limit" in error_str:
            return "系统繁忙，请稍后再试 🙏"
        elif "timeout" in error_str.lower() or "timed out" in error_str.lower():
            return "响应超时，请稍后再试"
        elif "500" in error_str or "502" in error_str or "503" in error_str:
            return "服务暂时不可用，请稍后再试"
        else:
            return "抱歉，处理图片时出现问题，请稍后再试"


@app.route("/webhook", methods=["POST"])
def webhook():
    """飞书事件回调"""
    data = request.json
    
    # URL 验证（首次配置时）
    if "challenge" in data:
        return jsonify({"challenge": data["challenge"]})
    
    # 处理消息事件
    header = data.get("header", {})
    event = data.get("event", {})
    
    # 消息去重：检查 event_id 是否已处理
    event_id = header.get("event_id", "")
    if message_dedup.is_duplicate(event_id):
        print(f"跳过重复事件: {event_id}")
        return jsonify({"code": 0})
    
    print(f"收到飞书事件: {json.dumps(data, ensure_ascii=False)}")
    
    # 验证 token (从 header 中获取)
    received_token = header.get("token")
    if received_token != FEISHU_VERIFICATION_TOKEN:
        print(f"Token 不匹配: {received_token} vs {FEISHU_VERIFICATION_TOKEN}")
        return jsonify({"error": "invalid token"}), 403
    
    if header.get("event_type") == "im.message.receive_v1":
        message = event.get("message", {})
        chat_id = message.get("chat_id")
        # 获取发送者信息
        sender = event.get("sender", {})
        sender_type = sender.get("sender_type", "user")  # "user" 或 "app"
        sender_info = sender.get("sender_id", {})
        sender_id = sender_info.get("user_id") or sender_info.get("open_id") or "unknown"
        msg_type = message.get("message_type")
        
        # 标记是否是 bot 消息
        is_bot_message = (sender_type == "app")
        if is_bot_message:
            print(f"📤 收到 Bot 消息: sender_id={sender_id}")
        
        # 保存用户的 chat_id (用于后续通知，如果是 admin 的话)
        open_id = sender_info.get("open_id", "")
        if not is_bot_message:
            save_user_chat_id(open_id, chat_id)
        
        # 白名单检查 - 但允许权限申请
        if not is_user_allowed(sender_info, sender_type):
            print(f"用户 {sender_id} 不在白名单中")
            
            chat_type = message.get("chat_type")  # p2p 或 group
            is_group_chat = (chat_type == "group")
            
            # 检查是否是权限申请相关消息
            if msg_type == "text":
                content = json.loads(message.get("content", "{}"))
                text = content.get("text", "").strip()
                
                # 移除 @mention 文本（修复：申请权限时也需要清理）
                import re
                text = re.sub(r'@\S+\s*', '', text).strip()
                
                # 获取用户名
                # 尝试从消息中获取发送者名字
                sender_name = sender_id  # 默认用 ID
                
                # 处理权限申请
                reply = handle_permission_request(text, open_id, chat_id, sender_name)
                if reply:
                    if is_group_chat:
                        # 群聊：简短提示 + 私聊发详情
                        send_feishu_message(chat_id, "📩 请查看私聊消息")
                        send_private_message(open_id, reply)
                    else:
                        # 私聊：直接回复
                        send_feishu_message(chat_id, reply)
                    return jsonify({"code": 0})
                
                # 不是申请相关消息，提示用户申请
                if is_group_chat:
                    send_feishu_message(chat_id, "你还没有使用权限，请私聊我发送「申请权限」。")
                else:
                    send_feishu_message(chat_id, "你还没有使用权限。\n\n发送「申请权限」来申请使用机器人。")
            
            return jsonify({"code": 0})
        
        # 处理文本消息
        if msg_type == "text":
            content = json.loads(message.get("content", "{}"))
            text = content.get("text", "")
            chat_type = message.get("chat_type")  # p2p 或 group
            
            # 群聊中需要 @机器人 才响应
            # 因为用了 group_at_msg 权限，收到群消息就意味着被 @ 了
            mentions = message.get("mentions", [])
            is_mentioned = len(mentions) > 0  # 有 mentions 就是被 @ 了
            
            # 移除 @mention 文本
            import re
            text = re.sub(r'@\S+\s*', '', text).strip()
            
            # 检查是否是验证码 (私聊时优先处理)
            if chat_type == "p2p" and is_verification_code(text):
                open_id = sender_info.get("open_id", "")
                print(f"收到验证码: {text}, open_id: {open_id}, chat_id: {chat_id}")
                result = verify_code_with_admin(text, open_id, chat_id)  # 传入 chat_id 用于后续通知
                if result.get("found"):
                    send_feishu_message(chat_id, f"✅ {result.get('message', '验证成功！')}")
                else:
                    send_feishu_message(chat_id, f"❌ {result.get('message', '验证码无效或已过期')}")
                return jsonify({"code": 0})
            
            # 私聊直接响应，群聊需要被 @ 或包含关键词
            should_respond = (chat_type == "p2p") or is_mentioned or ("clawdbot" in text.lower())
            
            if text.strip() and should_respond:
                print(f"用户 {sender_id}: {text}")
                
                # 检查是否在请求读取文件（群聊中的文件回看功能）
                if chat_type == "group" and is_file_request(text):
                    print(f"检测到文件读取请求: {text}")
                    send_feishu_message(chat_id, "🔍 正在查找最近的文件...")
                    
                    file_info = find_recent_file(chat_id)
                    if file_info.get("found"):
                        file_name = file_info["file_name"]
                        file_key = file_info["file_key"]
                        message_id = file_info["message_id"]
                        
                        # 检查是否是支持的文件类型
                        if not is_supported_file(file_name):
                            send_feishu_message(chat_id, f"找到文件 {file_name}，但不是支持的类型。\n支持: Excel, PDF, Word, PPT, Markdown, 文本文件, 视频, 音频")
                            return jsonify({"code": 0})
                        
                        # 根据文件类型选择图标
                        if is_excel_file(file_name):
                            icon = "📊"
                            file_type_name = "Excel"
                        elif is_pdf_file(file_name):
                            icon = "📄"
                            file_type_name = "PDF"
                        elif is_ppt_file(file_name):
                            icon = "📽️"
                            file_type_name = "PPT"
                        elif is_md_file(file_name):
                            icon = "📑"
                            file_type_name = "Markdown"
                        elif is_text_file(file_name):
                            icon = "📃"
                            file_type_name = "文本"
                        elif is_video_file(file_name):
                            icon = "🎬"
                            file_type_name = "视频"
                        elif is_audio_file(file_name):
                            icon = "🎵"
                            file_type_name = "音频"
                        else:
                            icon = "📝"
                            file_type_name = "Word"
                        
                        send_feishu_message(chat_id, f"{icon} 找到 {file_type_name} 文件: {file_name}，正在处理...")
                        
                        # 下载文件
                        temp_path = download_feishu_file(message_id, file_key, file_name)
                        if not temp_path:
                            send_feishu_message(chat_id, "❌ 文件下载失败，请稍后重试")
                            return jsonify({"code": 0})
                        
                        # 解析文档
                        doc_content, doc_type = parse_document_file(temp_path, file_name)
                        
                        # 构造发给 Clawdbot 的消息
                        prompt = f"用户发送了一个 {doc_type}: {file_name}\n\n文件内容如下:\n```\n{doc_content}\n```\n\n请分析这个文档内容，提供有用的见解或回答用户可能的问题。"
                        
                        # 调用 Clawdbot
                        open_id = sender_info.get("open_id", "")
                        reply = ask_clawdbot(prompt, sender_id, open_id, chat_id)
                        
                        if reply and reply.strip() != "No response from Clawdbot.":
                            send_feishu_message(chat_id, reply)
                            print(f"{doc_type} 分析回复: {reply[:100]}...")
                        else:
                            print("Clawdbot 无回复，保持静默")
                        
                        return jsonify({"code": 0})
                    else:
                        send_feishu_message(chat_id, "❌ 没有找到最近的文件。请先发送文件，然后 @我 说「读文件」。")
                        return jsonify({"code": 0})
                
                # 调用 Clawdbot (传入 open_id 和 chat_id)
                open_id = sender_info.get("open_id", "")
                reply = ask_clawdbot(text, sender_id, open_id, chat_id)
                
                # 过滤掉"无响应"消息，保持静默
                if reply and reply.strip() != "No response from Clawdbot.":
                    send_feishu_message(chat_id, reply)
                    print(f"回复: {reply[:100]}...")
                else:
                    print("Clawdbot 无回复，保持静默")
        
        # 处理富文本消息 (post) - 可能包含图片
        elif msg_type == "post":
            content = json.loads(message.get("content", "{}"))
            message_id = message.get("message_id", "")
            chat_type = message.get("chat_type")
            
            # 解析 post 内容，提取图片和文本
            post_content = content.get("content", [])
            image_keys = []
            text_parts = []
            
            for paragraph in post_content:
                for element in paragraph:
                    tag = element.get("tag", "")
                    if tag == "img":
                        image_keys.append(element.get("image_key", ""))
                    elif tag == "text":
                        text_parts.append(element.get("text", ""))
            
            # 提取的文本（去除空白）
            extracted_text = " ".join(text_parts).strip()
            
            print(f"收到富文本: images={len(image_keys)}, text='{extracted_text[:50]}...'")
            
            # 检查是否被 @ 或私聊
            mentions = message.get("mentions", [])
            is_mentioned = len(mentions) > 0
            should_respond = (chat_type == "p2p") or is_mentioned
            
            if not should_respond:
                print("群聊富文本未被 @，忽略")
                return jsonify({"code": 0})
            
            # 如果有图片，处理图片
            if image_keys:
                image_key = image_keys[0]  # 取第一张图片
                send_feishu_message(chat_id, "🖼️ 收到图片，正在处理...")
                
                # 下载图片
                temp_path = download_feishu_image(message_id, image_key)
                if not temp_path:
                    send_feishu_message(chat_id, "❌ 图片下载失败，请稍后重试")
                    return jsonify({"code": 0})
                
                # 发送给 Clawdbot（带上用户的文字描述）
                open_id = sender_info.get("open_id", "")
                prompt = extracted_text if extracted_text else "用户发送了一张图片，请分析这张图片。"
                reply = ask_clawdbot_with_image(temp_path, sender_id, open_id, chat_id, prompt)
                
                # 清理临时文件
                try:
                    os.remove(temp_path)
                except:
                    pass
                
                if reply and reply.strip() != "No response from Clawdbot.":
                    send_feishu_message(chat_id, reply)
                    print(f"图片分析回复: {reply[:100]}...")
                else:
                    print("Clawdbot 无回复，保持静默")
            else:
                # 没有图片，当作普通文本处理
                if extracted_text:
                    open_id = sender_info.get("open_id", "")
                    reply = ask_clawdbot(extracted_text, sender_id, open_id, chat_id)
                    
                    if reply and reply.strip() != "No response from Clawdbot.":
                        send_feishu_message(chat_id, reply)
                        print(f"回复: {reply[:100]}...")
                    else:
                        print("Clawdbot 无回复，保持静默")
        
        # 处理图片消息
        elif msg_type == "image":
            content = json.loads(message.get("content", "{}"))
            image_key = content.get("image_key", "")
            message_id = message.get("message_id", "")
            chat_type = message.get("chat_type")
            
            print(f"收到图片: image_key={image_key}, message_id={message_id}")
            
            # 私聊直接处理，群聊需要被 @ 才处理
            mentions = message.get("mentions", [])
            is_mentioned = len(mentions) > 0
            should_respond = (chat_type == "p2p") or is_mentioned
            
            if not should_respond:
                print("群聊图片未被 @，忽略")
                return jsonify({"code": 0})
            
            send_feishu_message(chat_id, "🖼️ 收到图片，正在处理...")
            
            # 下载图片
            temp_path = download_feishu_image(message_id, image_key)
            if not temp_path:
                send_feishu_message(chat_id, "❌ 图片下载失败，请稍后重试")
                return jsonify({"code": 0})
            
            # 发送给 Clawdbot（作为图片附件）
            open_id = sender_info.get("open_id", "")
            reply = ask_clawdbot_with_image(temp_path, sender_id, open_id, chat_id)
            
            # 清理临时文件
            try:
                os.remove(temp_path)
            except:
                pass
            
            if reply and reply.strip() != "No response from Clawdbot.":
                send_feishu_message(chat_id, reply)
                print(f"图片分析回复: {reply[:100]}...")
            else:
                print("Clawdbot 无回复，保持静默")
        
        # 处理文件消息 (Excel, PDF, Word)
        elif msg_type == "file":
            content = json.loads(message.get("content", "{}"))
            file_key = content.get("file_key", "")
            file_name = content.get("file_name", "")
            message_id = message.get("message_id", "")
            chat_type = message.get("chat_type")
            
            print(f"收到文件: {file_name}, file_key: {file_key}")
            
            # 检查是否是支持的文件类型
            if not is_supported_file(file_name):
                send_feishu_message(chat_id, f"支持的文件类型: Excel, PDF, Word, PPT, Markdown, 文本文件, 视频, 音频\n收到的是: {file_name}")
                return jsonify({"code": 0})
            
            # 私聊直接处理，群聊中支持的文件也直接处理（因为飞书群聊发文件不能同时 @）
            mentions = message.get("mentions", [])
            is_mentioned = len(mentions) > 0
            is_supported = is_supported_file(file_name)
            should_respond = (chat_type == "p2p") or is_mentioned or is_supported
            
            if not should_respond:
                print(f"群聊文件未被 @，忽略")
                return jsonify({"code": 0})
            
            # 根据文件类型选择提示图标
            if is_excel_file(file_name):
                icon = "📊"
                file_type_name = "Excel"
            elif is_pdf_file(file_name):
                icon = "📄"
                file_type_name = "PDF"
            elif is_ppt_file(file_name):
                icon = "📽️"
                file_type_name = "PPT"
            elif is_md_file(file_name):
                icon = "📑"
                file_type_name = "Markdown"
            elif is_text_file(file_name):
                icon = "📃"
                file_type_name = "文本"
            elif is_video_file(file_name):
                icon = "🎬"
                file_type_name = "视频"
            elif is_audio_file(file_name):
                icon = "🎵"
                file_type_name = "音频"
            else:
                icon = "📝"
                file_type_name = "Word"
            
            # 发送处理中提示
            send_feishu_message(chat_id, f"{icon} 正在处理 {file_type_name} 文件: {file_name}...")
            
            # 下载文件
            temp_path = download_feishu_file(message_id, file_key, file_name)
            if not temp_path:
                send_feishu_message(chat_id, "❌ 文件下载失败，请稍后重试")
                return jsonify({"code": 0})
            
            # 解析文档
            doc_content, doc_type = parse_document_file(temp_path, file_name)
            
            # 构造发给 Clawdbot 的消息
            prompt = f"用户发送了一个 {doc_type}: {file_name}\n\n文件内容如下:\n```\n{doc_content}\n```\n\n请分析这个文档内容，提供有用的见解或回答用户可能的问题。"
            
            # 调用 Clawdbot
            open_id = sender_info.get("open_id", "")
            reply = ask_clawdbot(prompt, sender_id, open_id, chat_id)
            
            # 回复（过滤无响应消息）
            if reply and reply.strip() != "No response from Clawdbot.":
                send_feishu_message(chat_id, reply)
                print(f"{doc_type} 分析回复: {reply[:100]}...")
            else:
                print("Clawdbot 无回复，保持静默")
        
        # 处理媒体消息 (视频、音频)
        elif msg_type == "media":
            content = json.loads(message.get("content", "{}"))
            file_key = content.get("file_key", "")
            file_name = content.get("file_name", "media_file")
            message_id = message.get("message_id", "")
            chat_type = message.get("chat_type")
            
            # 尝试从 image_key 获取（某些情况下飞书用 image_key）
            if not file_key:
                file_key = content.get("image_key", "")
            
            print(f"收到媒体文件: {file_name}, file_key: {file_key}")
            
            # 确定文件类型
            if is_video_file(file_name):
                icon = "🎬"
                file_type_name = "视频"
            elif is_audio_file(file_name):
                icon = "🎵"
                file_type_name = "音频"
            else:
                # 默认当视频处理
                icon = "🎬"
                file_type_name = "媒体"
                if not file_name.endswith(('.mp4', '.mov', '.avi', '.mp3', '.wav')):
                    file_name = file_name + ".mp4"  # 添加默认扩展名
            
            # 私聊直接处理，群聊也处理（媒体文件不能同时 @）
            chat_type = message.get("chat_type")
            mentions = message.get("mentions", [])
            is_mentioned = len(mentions) > 0
            should_respond = (chat_type == "p2p") or is_mentioned or True  # 媒体文件总是处理
            
            if not should_respond:
                print(f"群聊媒体未处理")
                return jsonify({"code": 0})
            
            # 发送处理中提示
            send_feishu_message(chat_id, f"{icon} 收到 {file_type_name} 文件: {file_name}，正在处理...")
            
            # 下载媒体文件
            temp_path = download_feishu_file(message_id, file_key, file_name)
            if not temp_path:
                send_feishu_message(chat_id, "❌ 媒体文件下载失败，请稍后重试")
                return jsonify({"code": 0})
            
            # 解析媒体文件（返回文件信息）
            doc_content, doc_type = parse_document_file(temp_path, file_name)
            
            # 构造发给 Clawdbot 的消息
            prompt = f"用户发送了一个 {doc_type}: {file_name}\n\n{doc_content}"
            
            # 调用 Clawdbot
            open_id = sender_info.get("open_id", "")
            reply = ask_clawdbot(prompt, sender_id, open_id, chat_id)
            
            # 回复
            if reply and reply.strip() != "No response from Clawdbot.":
                send_feishu_message(chat_id, reply)
                print(f"{doc_type} 处理回复: {reply[:100]}...")
            else:
                print("Clawdbot 无回复，保持静默")
    
    return jsonify({"code": 0})


@app.route("/health", methods=["GET"])
def health():
    """健康检查"""
    return jsonify({"status": "ok"})


@app.route("/notify", methods=["POST"])
def notify_user():
    """发送通知给用户 (供 bot-admin 调用)"""
    data = request.json
    chat_id = data.get("chat_id")
    open_id = data.get("open_id")  # 支持私聊
    message = data.get("message")
    
    if not message:
        return jsonify({"error": "缺少 message 参数"}), 400
    
    if not chat_id and not open_id:
        return jsonify({"error": "需要 chat_id 或 open_id"}), 400
    
    if open_id:
        # 私聊
        result = send_private_message(open_id, message)
    else:
        # 群聊
        result = send_feishu_message(chat_id, message)
    
    return jsonify({"success": True, "result": result})


@app.route("/send_file", methods=["POST"])
def send_file_endpoint():
    """发送文件到飞书聊天 (供 Clawdbot 调用)
    
    请求体:
    {
        "chat_id": "oc_xxx",      # 群聊ID (与 open_id 二选一)
        "open_id": "ou_xxx",      # 用户ID，用于私聊发送 (与 chat_id 二选一)
        "file_path": "/path/to/file",  # 文件路径
        "file_name": "xxx.xlsx",  # 可选，文件名
        "caption": "说明文字"     # 可选，发送文件前的说明
    }
    """
    data = request.json
    chat_id = data.get("chat_id")
    open_id = data.get("open_id")
    file_path = data.get("file_path")
    file_name = data.get("file_name")
    caption = data.get("caption")
    
    if not chat_id and not open_id:
        return jsonify({"success": False, "error": "需要 chat_id 或 open_id 参数"}), 400
    if not file_path:
        return jsonify({"success": False, "error": "缺少 file_path 参数"}), 400
    
    if not os.path.exists(file_path):
        return jsonify({"success": False, "error": f"文件不存在: {file_path}"}), 400
    
    result = send_file_to_chat(chat_id=chat_id, file_path=file_path, file_name=file_name, caption=caption, open_id=open_id)
    
    if result.get("success"):
        return jsonify(result)
    else:
        return jsonify(result), 500


@app.route("/card_callback", methods=["POST"])
def card_callback():
    """飞书消息卡片回调"""
    # 打印所有请求信息用于调试
    print(f"=== 卡片回调请求 ===")
    print(f"Headers: {dict(request.headers)}")
    print(f"Raw data: {request.get_data(as_text=True)}")
    
    data = request.json or {}
    print(f"收到卡片回调: {json.dumps(data, ensure_ascii=False)}")
    
    # URL 验证 (飞书首次配置时会发送 challenge)
    if "challenge" in data:
        return jsonify({"challenge": data["challenge"]})
    
    # 处理卡片按钮点击
    try:
        # 飞书卡片回调的 action 在 event 里面
        action = data.get("event", {}).get("action", {})
        action_value = action.get("value", "{}")
        
        print(f"卡片回调 action: {action}")
        
        # 解析按钮值
        if isinstance(action_value, str):
            btn_data = json.loads(action_value)
        else:
            btn_data = action_value
        
        action_type = btn_data.get("action")
        applicant_open_id = btn_data.get("open_id")
        applicant_name = btn_data.get("name")
        applicant_chat_id = btn_data.get("chat_id")
        selected_role = btn_data.get("role", "user")  # 从按钮值获取角色
        
        print(f"解析结果: action={action_type}, name={applicant_name}, role={selected_role}")
        
        if action_type == "approve":
            # 批准申请
            add_to_whitelist(applicant_open_id, applicant_name)
            add_to_permissions(applicant_open_id, applicant_name, selected_role)
            
            # 从待审批列表移除
            pending = load_pending()
            pending.get("pending", {}).pop(applicant_open_id, None)
            save_pending(pending)
            
            # 通知申请人
            send_feishu_message(applicant_chat_id, 
                f"🎉 你的权限申请已批准！\n\n角色: {selected_role} ({ROLE_DESCRIPTIONS.get(selected_role, '')})\n\n现在可以开始和我对话了。")
            
            # 返回更新后的卡片
            return jsonify({
                "toast": {"type": "success", "content": f"已批准 {applicant_name}，角色: {selected_role}"},
                "card": {
                    "config": {"wide_screen_mode": True},
                    "header": {"title": {"tag": "plain_text", "content": "✅ 已批准"}, "template": "green"},
                    "elements": [
                        {"tag": "div", "text": {"tag": "lark_md", "content": f"**{applicant_name}** 已获得 **{selected_role}** 权限"}}
                    ]
                }
            })
        
        elif action_type == "reject":
            # 拒绝申请
            pending = load_pending()
            pending.get("pending", {}).pop(applicant_open_id, None)
            save_pending(pending)
            
            # 通知申请人
            send_feishu_message(applicant_chat_id, 
                "❌ 你的权限申请被拒绝了。\n\n如有疑问请联系管理员。")
            
            # 返回更新后的卡片
            return jsonify({
                "toast": {"type": "info", "content": f"已拒绝 {applicant_name}"},
                "card": {
                    "config": {"wide_screen_mode": True},
                    "header": {"title": {"tag": "plain_text", "content": "❌ 已拒绝"}, "template": "red"},
                    "elements": [
                        {"tag": "div", "text": {"tag": "lark_md", "content": f"**{applicant_name}** 的申请已拒绝"}}
                    ]
                }
            })
    
    except Exception as e:
        print(f"处理卡片回调出错: {e}")
        import traceback
        traceback.print_exc()
    
    return jsonify({"success": True})


if __name__ == "__main__":
    print("飞书桥接服务启动...")
    print(f"Clawdbot: {CLAWDBOT_URL}")
    print("Webhook URL: http://YOUR_SERVER:5001/webhook")
    app.run(host="0.0.0.0", port=5001)
